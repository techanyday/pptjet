from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from openai import OpenAI
import os
import json
import requests
from typing import List, Dict, Optional

class PPTGenerator:
    def __init__(self):
        self.openai_api_key = os.environ.get('OPENAI_API_KEY')
        if not self.openai_api_key:
            raise ValueError("OpenAI API key not found in environment variables")
        self.client = OpenAI(api_key=self.openai_api_key)
        self.pexels_api_key = os.environ.get('PEXELS_API_KEY')
        
        # Define available template styles
        self.TEMPLATE_STYLES = {
            "Aesthetic": "Aesthetic.pptx",
            "Minimalist": "Minimalist.pptx",
            "Professional": "Professional.pptx",
            "Vintage": "Vintage.pptx"
        }
        
        # Define font settings (fallback for missing template styles)
        self.FONTS = {
            'title': {
                'name': 'Calibri',
                'size': Pt(44),
                'bold': True,
                'color': RGBColor(33, 33, 33)
            },
            'subtitle': {
                'name': 'Calibri',
                'size': Pt(32),
                'bold': False,
                'color': RGBColor(0, 123, 255)
            },
            'body': {
                'name': 'Calibri',
                'size': Pt(18),
                'bold': False,
                'color': RGBColor(33, 37, 41)
            }
        }

    def get_template_path(self, style: str) -> str:
        """Get the path to the selected template style"""
        filename = self.TEMPLATE_STYLES.get(style, "Professional.pptx")  # Default to Professional if style not found
        return os.path.join("app", "static", "presentations", "custom_styles", filename)

    def _remove_all_slides(self, prs: Presentation) -> None:
        """Remove all existing slides from the presentation while preserving the template"""
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)  # Create a list from iterator
        for slide_id in slides:
            xml_slides.remove(slide_id)

    def _apply_text_style(self, shape, style_type='body'):
        """Apply text style to a shape (only used if template styles are missing)"""
        if not shape.has_text_frame:
            return

        style = self.FONTS[style_type]
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = style['name']
            paragraph.font.size = style['size']
            paragraph.font.bold = style['bold']
            paragraph.font.color.rgb = style['color']
            if style_type in ['title', 'subtitle']:
                paragraph.alignment = PP_ALIGN.CENTER

    def _get_layout_by_name(self, prs: Presentation, layout_names: List[str]) -> Optional[any]:
        """Find a slide layout by trying multiple possible names"""
        # First try exact matches
        for name in layout_names:
            for layout in prs.slide_layouts:
                if layout.name.lower() == name.lower():
                    return layout
        
        # Then try partial matches
        for name in layout_names:
            for layout in prs.slide_layouts:
                if name.lower() in layout.name.lower():
                    return layout
        
        return None

    def _get_title_layout(self, prs: Presentation) -> any:
        """Get the title slide layout"""
        # Common names for title layouts
        title_layouts = [
            "Title Slide",
            "Title",
            "Cover",
            "Opening",
            "Front Page"
        ]
        
        # Try to find by name first
        layout = self._get_layout_by_name(prs, title_layouts)
        if layout:
            return layout
            
        # Fallback: Look for layout with title placeholder but no body
        for layout in prs.slide_layouts:
            has_title = False
            has_body = False
            for shape in layout.placeholders:
                if shape.placeholder_format.idx == 0:  # Title
                    has_title = True
                if shape.placeholder_format.idx == 1:  # Body
                    has_body = True
            if has_title and not has_body:
                return layout
        
        # Last resort: use first layout
        return prs.slide_layouts[0]

    def _get_content_layout(self, prs: Presentation) -> any:
        """Get the content slide layout"""
        # Common names for content layouts
        content_layouts = [
            "Title and Content",
            "Section",
            "Content",
            "Text",
            "Body",
            "Bullet",
            "Title and Text"
        ]
        
        # Try to find by name first
        layout = self._get_layout_by_name(prs, content_layouts)
        if layout:
            return layout
            
        # Fallback: Look for layout with both title and body placeholders
        for layout in prs.slide_layouts:
            has_title = False
            has_body = False
            for shape in layout.placeholders:
                if shape.placeholder_format.idx == 0:  # Title
                    has_title = True
                if shape.placeholder_format.idx == 1:  # Body
                    has_body = True
            if has_title and has_body:
                return layout
        
        # Last resort: use second layout if available, otherwise first
        return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    def get_named_layout(self, prs, preferred_names):
        """Find a layout by name from a list of preferred names"""
        for layout in prs.slide_layouts:
            if layout.name in preferred_names:
                return layout
        return prs.slide_layouts[0]  # fallback

    def _add_title_slide(self, prs: Presentation, title: str, presenter: str):
        """Add title slide"""
        preferred_names = ["Title Slide", "Title", "Opening Layout"]
        layout = self.get_named_layout(prs, preferred_names)
        slide = prs.slides.add_slide(layout)
        
        # Find and populate title placeholder
        title_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:  # Title
                title_placeholder = shape
                break
        
        if title_placeholder:
            title_placeholder.text = title
        
        # Find and populate subtitle/presenter placeholder
        subtitle = None
        for shape in slide.placeholders:
            # Look for subtitle or text placeholder
            if shape.placeholder_format.idx in [1, 2]:  # Common indices for subtitles
                subtitle = shape
                break
        
        if subtitle:
            subtitle.text = f"Presented by {presenter}"

    def _add_content_slide(self, prs: Presentation, title: str, content: str):
        """Add content slide"""
        preferred_names = ["Title and Content", "Title and Body", "Content", "Text and Content"]
        layout = self.get_named_layout(prs, preferred_names)
        slide = prs.slides.add_slide(layout)
        
        # Find and populate title placeholder
        title_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:  # Title
                title_placeholder = shape
                break
        
        if title_placeholder:
            title_placeholder.text = title
        
        # Find content placeholder - try multiple indices
        content_placeholder = None
        for shape in slide.placeholders:
            # Common indices for content/body placeholders
            if shape.placeholder_format.idx in [1, 2, 3]:
                content_placeholder = shape
                break
        
        # Add content if placeholder exists
        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.word_wrap = True
            
            # Add bullet points
            first = True
            for point in content.split('\n'):
                if first:
                    p = text_frame.paragraphs[0]
                    first = False
                else:
                    p = text_frame.add_paragraph()
                
                p.text = point.strip()
                p.level = 0
                p.space_before = Pt(12)
                p.space_after = Pt(12)
                p.line_spacing = 1.2

    def _debug_print_layouts(self, prs: Presentation):
        """Print all available layouts in the presentation for debugging"""
        print(f"\nPresentation Details:")
        print(f"Slide Masters: {len(prs.slide_masters)}")
        for i, master in enumerate(prs.slide_masters):
            print(f"\nMaster {i}:")
            print(f"  Background: {master.background}")
            print(f"  Slide Layouts: {len(master.slide_layouts)}")
            for j, layout in enumerate(master.slide_layouts):
                print(f"\n  Layout {j}: {layout.name}")
                print("    Placeholders:")
                for shape in layout.placeholders:
                    print(f"      - {shape.name} (idx: {shape.placeholder_format.idx})")
                print("    Theme:")
                if hasattr(layout, 'theme'):
                    print(f"      Font Theme: {layout.theme.font_scheme.name if layout.theme.font_scheme else 'None'}")
        print("\n")

    def get_relevant_image(self, query: str) -> Optional[str]:
        """Get relevant image URL from Pexels API"""
        if not self.pexels_api_key:
            return None
        
        headers = {'Authorization': self.pexels_api_key}
        response = requests.get(
            f'https://api.pexels.com/v1/search?query={query}&per_page=1',
            headers=headers
        )
        
        if response.status_code == 200:
            data = response.json()
            if data['photos']:
                return data['photos'][0]['src']['large']
        return None

    def generate_slide_content(self, prompt: str, num_slides: int) -> List[Dict]:
        """Generate slide content using GPT-3.5"""
        try:
            if not self.openai_api_key:
                raise ValueError("OpenAI API key not set")
            
            system_prompt = f"""Generate a PowerPoint presentation with exactly {num_slides} slides.
            Return your response as a valid JSON object with a 'slides' array. Each slide must include:
            1. title (string): A concise, engaging title
            2. content (array): A list of 3-5 bullet points

            Format example:
            {{
                "slides": [
                    {{
                        "title": "Introduction",
                        "content": [
                            "Key point 1",
                            "Key point 2",
                            "Key point 3"
                        ]
                    }},
                    ...
                ]
            }}

            Topic: {prompt}
            
            Remember:
            - Keep titles concise and clear
            - Each bullet point should be a complete thought
            - Ensure the content flows logically
            - Return ONLY the JSON object, no other text"""

            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                response_format={ "type": "json_object" }
            )
            
            # Parse the JSON response
            content = response.choices[0].message.content
            print(f"GPT Response: {content}")  # Debug print
            response_data = json.loads(content)
            
            if not isinstance(response_data, dict) or 'slides' not in response_data:
                raise ValueError("Invalid response format from GPT. Expected object with 'slides' array.")
            
            slides_data = response_data['slides']
            
            if not isinstance(slides_data, list):
                raise ValueError("Invalid 'slides' format. Expected array.")
            
            # Ensure we have the right number of slides and format the content
            slides = []
            for slide in slides_data[:num_slides]:
                if not isinstance(slide, dict) or 'title' not in slide or 'content' not in slide:
                    raise ValueError("Invalid slide format. Expected object with 'title' and 'content'.")
                
                if not isinstance(slide['content'], list):
                    raise ValueError("Invalid content format. Expected array of strings.")
                
                # Join content without bullet points - PowerPoint will add them automatically
                formatted_content = "\n".join(point.strip() for point in slide['content'])
                slides.append({
                    "title": slide['title'],
                    "content": formatted_content
                })
            
            return slides
        except json.JSONDecodeError as e:
            raise Exception(f"Error parsing GPT response: {str(e)}")
        except ValueError as e:
            raise Exception(f"Error in response format: {str(e)}")
        except Exception as e:
            raise Exception(f"Error generating content: {str(e)}")

    def _copy_slide_layout(self, source_layout, target_prs):
        """Copy a slide layout from template to new presentation"""
        # Get the slide master
        if len(target_prs.slide_masters) == 0:
            target_prs.slide_masters.add_slide_master()
        
        target_master = target_prs.slide_masters[0]
        
        # Copy theme information if available
        if hasattr(source_layout, 'theme'):
            # Copy font scheme
            if source_layout.theme.font_scheme:
                if hasattr(target_master, 'theme') and target_master.theme:
                    target_master.theme.font_scheme = source_layout.theme.font_scheme
        
        # Copy background if available
        if hasattr(source_layout, 'background'):
            for layout in target_master.slide_layouts:
                if hasattr(layout, 'background'):
                    layout.background = source_layout.background
        
        return target_master.slide_layouts[0]  # Return the first layout as fallback

    def create_presentation(self,
                        title: str,
                        presenter: str,
                        slides_content: List[Dict],
                        template: str = "Professional",
                        include_images: bool = False) -> str:
        """Create PowerPoint presentation using a selected template style"""
        # Load template
        template_path = self.get_template_path(template)
        try:
            template_prs = Presentation(template_path)
            # Debug: Print available layouts
            self._debug_print_layouts(template_prs)
            
            # Create new presentation
            prs = Presentation()
            
            # Copy master slides and layouts from template
            if len(template_prs.slide_masters) > 0:
                template_master = template_prs.slide_masters[0]
                if len(prs.slide_masters) == 0:
                    prs.slide_masters.add_slide_master()
                target_master = prs.slide_masters[0]
                
                # Copy theme information
                if hasattr(template_master, 'theme') and template_master.theme:
                    target_master.theme = template_master.theme
                
                # Copy background
                if hasattr(template_master, 'background'):
                    target_master.background = template_master.background
                
                # Copy layouts
                for layout in template_master.slide_layouts:
                    self._copy_slide_layout(layout, prs)
            
        except Exception as e:
            print(f"Warning: Could not load template {template_path}. Using blank presentation. Error: {str(e)}")
            prs = Presentation()

        # Add slides
        self._add_title_slide(prs, title, presenter)
        
        for slide_content in slides_content:
            self._add_content_slide(prs, slide_content['title'], slide_content['content'])

            if include_images and self.pexels_api_key:
                image_url = self.get_relevant_image(slide_content['title'])
                if image_url:
                    # TODO: Implement image insertion here if needed
                    pass

        # Ensure output directory exists
        os.makedirs('generated', exist_ok=True)

        # Clean filename
        safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '_')).rstrip()
        filename = f"{safe_title.replace(' ', '_')}.pptx"
        output_path = os.path.join('generated', filename)

        prs.save(output_path)
        return output_path
